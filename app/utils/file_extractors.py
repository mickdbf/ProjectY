import fitz
from pptx import Presentation

def extract_text_from_pdf(path):
    text = ""
    with fitz.open(path) as doc:
        for page in doc:
            text += page.get_text()
    return text.strip()

def extract_text_from_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text).strip()
